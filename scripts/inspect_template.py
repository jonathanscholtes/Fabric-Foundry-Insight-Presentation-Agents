import re, sys
from pptx import Presentation

sys.stdout.reconfigure(encoding="utf-8", errors="replace")

prs = Presentation(r"c:\Users\joscholt\OneDrive - Microsoft\Documents\GitHub\Azure-Fabric-MBR-AI-Agents\data\templates\mbr_template.pptx")
print("Slides:", len(prs.slides))

# Collect all unique tags across the whole deck
all_tags_single  = set()
all_tags_double  = set()

for i, slide in enumerate(prs.slides, 1):
    tags_single = []
    tags_double = []
    texts = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full = "".join(r.text for r in para.runs)
            if full.strip():
                texts.append(full.strip()[:100])
            tags_single.extend(re.findall(r"\{(\w+)\}", full))
            tags_double.extend(re.findall(r"\{\{(\w+)\}\}", full))
    all_tags_single.update(tags_single)
    all_tags_double.update(tags_double)
    print(f"\nSlide {i}:")
    if tags_single:  print("  {single} tags:", tags_single)
    if tags_double:  print("  {{double}} tags:", tags_double)
    if not tags_single and not tags_double:
        print("  Tags: (none)")
    for t in texts[:10]:
        print("  |", t)

print("\n=== Summary ===")
print("{single} tags found:", sorted(all_tags_single))
print("{{double}} tags found:", sorted(all_tags_double))
