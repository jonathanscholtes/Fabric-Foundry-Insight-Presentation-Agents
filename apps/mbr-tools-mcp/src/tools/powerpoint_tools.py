"""MCP tools: get_template_slides, fill_mbr_template, get_mbr_deck_url."""

from __future__ import annotations

import json
import logging
import os
import re
import subprocess
import tempfile
import uuid
from datetime import datetime, timezone, timedelta
from pathlib import Path

from azure.identity import ManagedIdentityCredential, DefaultAzureCredential
from azure.storage.blob import BlobServiceClient, generate_blob_sas, BlobSasPermissions
from azure.storage.blob import UserDelegationKey
from pptx import Presentation
from pptx.util import Pt

from ..config import settings

logger = logging.getLogger("mbr-tools-mcp.powerpoint")

SAS_EXPIRY_HOURS = 1


# ── Storage helpers ────────────────────────────────────────────────────────────

def _get_blob_client() -> BlobServiceClient:
    credential = (
        ManagedIdentityCredential(client_id=settings.AZURE_CLIENT_ID)
        if settings.AZURE_CLIENT_ID
        else DefaultAzureCredential()
    )
    return BlobServiceClient(account_url=settings.STORAGE_ACCOUNT_URL, credential=credential)


def _sas_url(blob_client: BlobServiceClient, container: str, blob_path: str) -> str:
    """Generate a short-lived user-delegation SAS URL."""
    expiry = datetime.now(timezone.utc) + timedelta(hours=SAS_EXPIRY_HOURS)
    start  = datetime.now(timezone.utc) - timedelta(minutes=5)

    key: UserDelegationKey = blob_client.get_user_delegation_key(start, expiry)
    account_name = blob_client.account_name

    token = generate_blob_sas(
        account_name=account_name,
        container_name=container,
        blob_name=blob_path,
        user_delegation_key=key,
        permission=BlobSasPermissions(read=True),
        expiry=expiry,
    )
    return f"{settings.STORAGE_ACCOUNT_URL}/{container}/{blob_path}?{token}"


def _upload_bytes(
    blob_client: BlobServiceClient, container: str, blob_path: str,
    data: bytes, content_type: str = "application/octet-stream",
) -> None:
    blob = blob_client.get_blob_client(container=container, blob=blob_path)
    blob.upload_blob(data, overwrite=True, content_settings={"content_type": content_type})


def _download_bytes(blob_client: BlobServiceClient, container: str, blob_path: str) -> bytes:
    blob = blob_client.get_blob_client(container=container, blob=blob_path)
    return blob.download_blob().readall()


# ── KPI formatting ─────────────────────────────────────────────────────────────

def _fmt_revenue(v: float) -> str:
    return f"${v / 1_000_000:.2f}M"

def _fmt_miles(v: float) -> str:
    return f"{v / 1_000_000:.2f}M mi"

def _fmt_pct(v: float) -> str:
    return f"{v:.1f}%"

def _fmt_cpm(v: float) -> str:
    return f"${v:.2f}"

def _fmt_delta_pct(v: float) -> str:
    arrow = "↑" if v >= 0 else "↓"
    return f"{arrow} {abs(v):.1f}%"

def _fmt_delta_pp(v: float) -> str:
    arrow = "↑" if v >= 0 else "↓"
    return f"{arrow} {abs(v):.1f}pp"


# ── PowerPoint placeholder filling ────────────────────────────────────────────

def _replace_in_shape(shape, replacements: dict[str, str]) -> None:
    """Replace {{tag}} placeholders in all text runs of a shape."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            for tag, value in replacements.items():
                if f"{{{{{tag}}}}}" in run.text:
                    run.text = run.text.replace(f"{{{{{tag}}}}}", value)


def _fill_slide(slide, replacements: dict[str, str]) -> None:
    for shape in slide.shapes:
        _replace_in_shape(shape, replacements)


def _build_kpi_replacements(kpis: dict) -> dict[str, str]:
    r = kpis
    return {
        "total_revenue":          _fmt_revenue(r["total_revenue"]["value"]),
        "total_revenue_delta":    _fmt_delta_pct(r["total_revenue"]["delta_pct"]),
        "total_miles":            _fmt_miles(r["total_miles"]["value"]),
        "total_miles_delta":      _fmt_delta_pct(r["total_miles"]["delta_pct"]),
        "empty_miles_pct":        _fmt_pct(r["empty_miles_pct"]["value"]),
        "empty_miles_delta":      _fmt_delta_pp(r["empty_miles_pct"]["delta_pp"]),
        "operating_margin_pct":   _fmt_pct(r["operating_margin_pct"]["value"]),
        "operating_margin_delta": _fmt_delta_pp(r["operating_margin_pct"]["delta_pp"]),
        "cost_per_mile":          _fmt_cpm(r["cost_per_mile"]["value"]),
        "cost_per_mile_delta":    _fmt_delta_pct(r["cost_per_mile"]["delta_pct"]),
        "on_time_delivery_pct":   _fmt_pct(r["on_time_delivery_pct"]["value"]),
        "on_time_delivery_delta": _fmt_delta_pp(r["on_time_delivery_pct"]["delta_pp"]),
    }


# ── LibreOffice thumbnail generation ──────────────────────────────────────────

def _generate_thumbnails(pptx_path: str, out_dir: str) -> list[str]:
    """Convert pptx slides to PNG files using LibreOffice headless.

    Returns sorted list of generated PNG paths.
    """
    result = subprocess.run(
        [
            "libreoffice", "--headless",
            "--convert-to", "png",
            "--outdir", out_dir,
            pptx_path,
        ],
        capture_output=True, text=True, timeout=120,
    )
    if result.returncode != 0:
        logger.error("LibreOffice conversion failed: %s", result.stderr)
        raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")

    stem = Path(pptx_path).stem
    pngs = sorted(Path(out_dir).glob(f"{stem}*.png"))
    return [str(p) for p in pngs]


# ── Tool implementations ───────────────────────────────────────────────────────

def register_powerpoint_tools(mcp) -> None:

    @mcp.tool()
    def get_template_slides(template_name: str) -> list[dict]:
        """Return pre-rendered slide thumbnails for the named template.

        Args:
            template_name: Template name, e.g. "mbr_template".

        Returns:
            List of {slide_number, title, thumbnail_url}.
        """
        client = _get_blob_client()
        prefix = f"templates/{template_name}/"
        container_client = client.get_container_client("thumbnails")

        slides: list[dict] = []
        slide_titles = [
            "Title",
            "Executive Summary",
            "Regional Performance",
            "Fleet Efficiency",
            "Driver Performance",
            "Bottom Line",
        ]

        for blob in container_client.list_blobs(name_starts_with=prefix):
            blob_name = blob.name
            filename = blob_name.split("/")[-1]  # e.g. "slide-01.png"
            m = re.match(r"slide-(\d+)\.png", filename)
            if not m:
                continue
            num = int(m.group(1))
            title = slide_titles[num - 1] if 1 <= num <= len(slide_titles) else f"Slide {num}"
            url = _sas_url(client, "thumbnails", blob_name)
            slides.append({"slide_number": num, "title": title, "thumbnail_url": url})

        slides.sort(key=lambda s: s["slide_number"])
        return slides

    @mcp.tool()
    def fill_mbr_template(
        region: str,
        period: str,
        kpis: dict,
        narratives: dict,
    ) -> dict:
        """Fill the MBR PowerPoint template with KPIs and narrative text blocks.

        Downloads mbr_template.pptx, fills all {{tag}} placeholders, uploads the
        completed deck, generates slide thumbnails with LibreOffice, and returns
        the deck ID, deck URL, and thumbnail URLs.

        Args:
            region:     Region name, e.g. "Southwest".
            period:     Period label, e.g. "May 2025".
            kpis:       Dict of KPI values and deltas from the Fabric Data Agent.
            narratives: Dict of text blocks: executive_summary, key_drivers,
                        cost_management, operational_efficiency, service_performance,
                        bottom_line.

        Returns:
            {deck_id, deck_url, thumbnail_urls}
        """
        client = _get_blob_client()
        deck_id = str(uuid.uuid4())[:8]
        period_slug = period.replace(" ", "")
        generated_date = datetime.now(timezone.utc).strftime("%B %d, %Y")

        kpi_tags = _build_kpi_replacements(kpis)

        # Slide-level replacements
        title_replacements = {
            "region": region,
            "period": period,
            "generated_date": generated_date,
        }
        exec_replacements = {**kpi_tags, "executive_summary": narratives.get("executive_summary", "")}
        regional_replacements = {"regional_performance_narrative": narratives.get("service_performance", "")}
        fleet_replacements = {"fleet_efficiency_narrative": (
            narratives.get("operational_efficiency", "") + "\n" + narratives.get("cost_management", "")
        )}
        driver_replacements = {"driver_performance_narrative": narratives.get("key_drivers", "")}
        bottom_replacements = {"bottom_line_narrative": narratives.get("bottom_line", "")}

        with tempfile.TemporaryDirectory() as tmpdir:
            # Download template
            template_bytes = _download_bytes(client, "templates", "mbr_template.pptx")
            template_path = os.path.join(tmpdir, "mbr_template.pptx")
            with open(template_path, "wb") as f:
                f.write(template_bytes)

            # Fill template
            prs = Presentation(template_path)
            slide_replacements = [
                title_replacements,
                exec_replacements,
                regional_replacements,
                fleet_replacements,
                driver_replacements,
                bottom_replacements,
            ]
            for idx, slide in enumerate(prs.slides):
                if idx < len(slide_replacements):
                    _fill_slide(slide, slide_replacements[idx])

            # Save filled deck
            deck_filename = f"{region}-{period_slug}-{deck_id}.pptx"
            deck_path = os.path.join(tmpdir, deck_filename)
            prs.save(deck_path)

            # Upload deck to Storage
            with open(deck_path, "rb") as f:
                deck_bytes = f.read()
            deck_blob_path = f"{region}-{period_slug}-{deck_id}.pptx"
            _upload_bytes(client, "decks", deck_blob_path, deck_bytes,
                          "application/vnd.openxmlformats-officedocument.presentationml.presentation")

            # Generate slide thumbnails
            thumb_dir = os.path.join(tmpdir, "thumbs")
            os.makedirs(thumb_dir, exist_ok=True)
            thumbnail_urls: list[str] = []

            try:
                png_paths = _generate_thumbnails(deck_path, thumb_dir)
                for i, png_path in enumerate(png_paths, start=1):
                    blob_path = f"decks/{deck_id}/slide-{i:02d}.png"
                    with open(png_path, "rb") as f:
                        png_bytes = f.read()
                    _upload_bytes(client, "thumbnails", blob_path, png_bytes, "image/png")
                    thumbnail_urls.append(_sas_url(client, "thumbnails", blob_path))
            except Exception as exc:
                logger.warning("Thumbnail generation failed, continuing without thumbnails: %s", exc)

            deck_url = _sas_url(client, "decks", deck_blob_path)

            # Save deck metadata blob
            metadata = {
                "deck_id":        deck_id,
                "period":         period,
                "region":         region,
                "generated_at":   datetime.now(timezone.utc).isoformat(),
                "deck_blob":      f"decks/{deck_blob_path}",
                "thumbnail_urls": thumbnail_urls,
            }
            _upload_bytes(
                client, "decks-metadata", f"{deck_id}.json",
                json.dumps(metadata).encode(), "application/json",
            )

            logger.info("fill_mbr_template complete: deck_id=%s region=%s period=%s", deck_id, region, period)
            return {"deck_id": deck_id, "deck_url": deck_url, "thumbnail_urls": thumbnail_urls}

    @mcp.tool()
    def get_mbr_deck_url(deck_id: str) -> str:
        """Return a short-lived SAS URL for downloading a completed MBR deck.

        Args:
            deck_id: The deck ID returned by fill_mbr_template.

        Returns:
            SAS URL string for the .pptx blob.
        """
        client = _get_blob_client()
        container_client = client.get_container_client("decks")
        for blob in container_client.list_blobs():
            if deck_id in blob.name:
                return _sas_url(client, "decks", blob.name)
        raise ValueError(f"No deck found for deck_id={deck_id!r}")
