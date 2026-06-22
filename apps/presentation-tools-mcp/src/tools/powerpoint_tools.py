"""MCP tools: get_template_slides, fill_presentation_template, get_deck_url."""

from __future__ import annotations

import json
import logging
import os
import re
import subprocess
import tempfile
import uuid
from datetime import datetime, timezone, timedelta
from pathlib import Path

from azure.identity import ManagedIdentityCredential, DefaultAzureCredential
from azure.storage.blob import BlobServiceClient, ContentSettings, generate_blob_sas, BlobSasPermissions
from azure.storage.blob import UserDelegationKey
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Pt

from ..config import settings

logger = logging.getLogger("presentation-tools.powerpoint")

SAS_EXPIRY_HOURS = 1


# ── Storage helpers ────────────────────────────────────────────────────────────

def _get_blob_client() -> BlobServiceClient:
    credential = (
        ManagedIdentityCredential(client_id=settings.AZURE_CLIENT_ID)
        if settings.AZURE_CLIENT_ID
        else DefaultAzureCredential()
    )
    return BlobServiceClient(account_url=settings.STORAGE_ACCOUNT_URL, credential=credential)


def _sas_url(blob_client: BlobServiceClient, container: str, blob_path: str) -> str:
    """Generate a short-lived user-delegation SAS URL."""
    expiry = datetime.now(timezone.utc) + timedelta(hours=SAS_EXPIRY_HOURS)
    start  = datetime.now(timezone.utc) - timedelta(minutes=5)

    key: UserDelegationKey = blob_client.get_user_delegation_key(start, expiry)
    account_name = blob_client.account_name

    token = generate_blob_sas(
        account_name=account_name,
        container_name=container,
        blob_name=blob_path,
        user_delegation_key=key,
        permission=BlobSasPermissions(read=True),
        expiry=expiry,
    )
    return f"{settings.STORAGE_ACCOUNT_URL}/{container}/{blob_path}?{token}"


def _upload_bytes(
    blob_client: BlobServiceClient, container: str, blob_path: str,
    data: bytes, content_type: str = "application/octet-stream",
) -> None:
    blob = blob_client.get_blob_client(container=container, blob=blob_path)
    blob.upload_blob(data, overwrite=True, content_settings=ContentSettings(content_type=content_type))


def _download_bytes(blob_client: BlobServiceClient, container: str, blob_path: str) -> bytes:
    blob = blob_client.get_blob_client(container=container, blob=blob_path)
    return blob.download_blob().readall()


# ── KPI formatting ─────────────────────────────────────────────────────────────

def _fmt_revenue(v: float) -> str:
    return f"${v / 1_000_000:.2f}M"

def _fmt_miles(v: float) -> str:
    return f"{v / 1_000_000:.2f}M mi"

def _fmt_pct(v: float) -> str:
    return f"{v:.1f}%"

def _fmt_cpm(v: float) -> str:
    return f"${v:.2f}"

def _fmt_delta_pct(v: float) -> str:
    arrow = "↑" if v >= 0 else "↓"
    return f"{arrow} {abs(v):.1f}%"

def _fmt_delta_pp(v: float) -> str:
    arrow = "↑" if v >= 0 else "↓"
    return f"{arrow} {abs(v):.1f}pp"


# ── Text splitting ─────────────────────────────────────────────────────────────

def _split_text(text: str, n: int) -> list[str]:
    """Split a narrative block into n parts by bullet lines, then sentences."""
    lines = [l.strip().lstrip("•-–*· ").strip() for l in text.splitlines() if l.strip()]
    if len(lines) >= n:
        return lines[:n]
    sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", text.strip()) if s.strip()]
    if len(sentences) >= n:
        return sentences[:n]
    return (lines + [""] * n)[:n]


# ── PowerPoint placeholder filling ────────────────────────────────────────────

def _replace_in_shape(shape, replacements: dict[str, str]) -> None:
    """Replace {tag} placeholders in all text runs of a shape."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            for tag, value in replacements.items():
                placeholder = f"{{{tag}}}"
                if placeholder in run.text:
                    run.text = run.text.replace(placeholder, value)


def _fill_slide(slide, replacements: dict[str, str]) -> None:
    for shape in slide.shapes:
        _replace_in_shape(shape, replacements)


def _build_kpi_tags(kpis: dict) -> dict[str, str]:
    """Map trucking KPI dict to the template's {tag} names."""
    r = kpis
    rev   = r["total_revenue"]
    miles = r["total_miles"]
    empty = r["empty_miles_pct"]
    margin = r["operating_margin_pct"]
    cpm   = r["cost_per_mile"]
    otd   = r["on_time_delivery_pct"]

    # load_factor is the inverse of empty miles pct
    load_factor       = 100.0 - empty["value"]
    load_factor_delta = -empty["delta_pp"]

    return {
        # Slide 2 — Executive Summary KPI boxes
        "total_revenue":          _fmt_revenue(rev["value"]),
        "revenue_delta":          _fmt_delta_pct(rev["delta_pct"]),
        "total_users":            _fmt_miles(miles["value"]),          # slot re-used for Total Miles
        "users_delta":            _fmt_delta_pct(miles["delta_pct"]),
        "gross_margin":           _fmt_pct(empty["value"]),            # slot re-used for Empty Miles %
        "gross_margin_delta":     _fmt_delta_pp(empty["delta_pp"]),
        "operating_margin":       _fmt_pct(margin["value"]),
        "operating_margin_delta": _fmt_delta_pp(margin["delta_pp"]),
        "cost_per_hour":          _fmt_cpm(cpm["value"]),              # slot re-used for Cost Per Mile
        "cost_per_hour_delta":    _fmt_delta_pct(cpm["delta_pct"]),
        "on_time_delivery":       _fmt_pct(otd["value"]),
        "on_time_delivery_delta": _fmt_delta_pp(otd["delta_pp"]),
        # Slide 3 — Revenue Performance
        "revenue_prior":          _fmt_revenue(r["revenue_prior"]["value"]) if r.get("revenue_prior") else "",
        # Slide 4 — Cost Management
        "total_cost":             _fmt_revenue(r["total_cost"]["value"]) if r.get("total_cost") else "",
        "total_cost_delta":       _fmt_delta_pct(r["total_cost"]["delta_pct"]) if r.get("total_cost") else "",
        # Slide 5 — Operational Efficiency
        "utilization_rate":       _fmt_pct(load_factor),
        "utilization_rate_delta": _fmt_delta_pp(load_factor_delta),
        "sla_compliance":         _fmt_pct(otd["value"]),
        "sla_compliance_delta":   _fmt_delta_pp(otd["delta_pp"]),
        "efficiency_score":       "",
    }


# ── Chart data population ──────────────────────────────────────────────────────

def _fill_charts(prs, analytics: dict | None, kpis: dict, region: str) -> None:
    """Replace data in the named chart shapes inserted into the template.

    Charts are matched by ``shape.name`` (case-sensitive): ``revenue_trend_chart``,
    ``cost_breakdown_chart``, ``sector_revenue_chart``. Any chart not present in the
    template is silently skipped — so this is safe to ship before the manual
    template update; it simply does nothing until the charts exist.

    The code only replaces chart *data*; styling stays whatever the template defines.
    """
    if not analytics:
        return

    builders: dict[str, CategoryChartData] = {}

    # Revenue trend — one series, 6 monthly points
    trend = analytics.get("revenue_trend")
    if trend:
        cd = CategoryChartData()
        cd.categories = [pt["period"] for pt in trend]
        cd.add_series("Revenue ($M)", [pt["revenue"] for pt in trend])
        builders["revenue_trend_chart"] = cd

    # Cost breakdown — one series, one category per cost bucket
    breakdown = analytics.get("cost_breakdown")
    if breakdown:
        cd = CategoryChartData()
        cd.categories = list(breakdown.keys())
        cd.add_series("Amount", list(breakdown.values()))
        builders["cost_breakdown_chart"] = cd

    # Sector revenue — one row per region from the kpis already available
    # (single-region decks today; the chart still renders the current region's bar).
    rev = kpis.get("total_revenue") if kpis else None
    if rev and rev.get("value") is not None:
        cd = CategoryChartData()
        cd.categories = [region]
        cd.add_series("Revenue ($M)", [round(rev["value"] / 1_000_000, 2)])
        builders["sector_revenue_chart"] = cd

    if not builders:
        return

    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_chart", False) and shape.name in builders:
                shape.chart.replace_data(builders[shape.name])
                logger.info("Filled chart shape: %s", shape.name)


# ── LibreOffice thumbnail generation ──────────────────────────────────────────

def _generate_thumbnails(pptx_path: str, out_dir: str) -> list[str]:
    """Convert pptx slides to PNG files using LibreOffice headless."""
    result = subprocess.run(
        [
            "libreoffice", "--headless",
            "--convert-to", "png",
            "--outdir", out_dir,
            pptx_path,
        ],
        capture_output=True, text=True, timeout=120,
    )
    if result.returncode != 0:
        logger.error("LibreOffice conversion failed: %s", result.stderr)
        raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")

    stem = Path(pptx_path).stem
    pngs = sorted(Path(out_dir).glob(f"{stem}*.png"))
    return [str(p) for p in pngs]


# ── Tool implementations ───────────────────────────────────────────────────────

def register_powerpoint_tools(mcp) -> None:

    @mcp.tool()
    def get_template_slides(template_name: str) -> list[dict]:
        """Return pre-rendered slide thumbnails for the named template.

        Args:
            template_name: Template name, e.g. "mbr_template".

        Returns:
            List of {slide_number, title, thumbnail_url}.
        """
        client = _get_blob_client()
        prefix = f"templates/{template_name}/"
        container_client = client.get_container_client("thumbnails")

        slide_titles = [
            "Title",
            "Executive Summary",
            "Revenue Performance",
            "Cost Management",
            "Operational Efficiency",
            "Sector Performance",
            "Key Drivers & Bottom Line",
            "Data Sources & Methodology",
        ]

        slides: list[dict] = []
        for blob in container_client.list_blobs(name_starts_with=prefix):
            blob_name = blob.name
            filename = blob_name.split("/")[-1]
            m = re.match(r"slide-(\d+)\.png", filename)
            if not m:
                continue
            num = int(m.group(1))
            title = slide_titles[num - 1] if 1 <= num <= len(slide_titles) else f"Slide {num}"
            url = _sas_url(client, "thumbnails", blob_name)
            slides.append({"slide_number": num, "title": title, "thumbnail_url": url})

        slides.sort(key=lambda s: s["slide_number"])
        return slides

    @mcp.tool()
    def fill_presentation_template(
        region: str,
        period: str,
        kpis: dict,
        narratives: dict,
        analytics: dict = None,
    ) -> dict:
        """Fill the MBR PowerPoint template with KPIs and narrative text blocks.

        Downloads mbr_template.pptx, fills all {tag} placeholders, populates any
        chart shapes when analytics data is supplied, uploads the completed deck,
        generates slide thumbnails with LibreOffice, and returns the deck ID,
        deck URL, and thumbnail URLs.

        Args:
            region:     Region name, e.g. "North".
            period:     Period label, e.g. "May 2025".
            kpis:       Dict of KPI values and deltas from the Fabric Data Agent.
            narratives: Dict of text blocks: executive_summary, key_drivers,
                        cost_management, operational_efficiency, service_performance,
                        bottom_line, call_to_action.
            analytics:  Optional dict with chart data:
                        {"revenue_trend": [{"period", "revenue"}, ...],
                         "cost_breakdown": {"Fuel", "Labor", "Maintenance", "Other"}}.
                        When omitted, charts are left as-is.

        Returns:
            {deck_id, deck_url, thumbnail_urls}
        """
        client = _get_blob_client()
        deck_id        = str(uuid.uuid4())[:8]
        period_slug    = period.replace(" ", "")
        generated_date = datetime.now(timezone.utc).strftime("%B %d, %Y")

        kpi_tags = _build_kpi_tags(kpis)

        # Split narrative blocks for slides that have multiple placeholder slots
        drivers      = _split_text(narratives.get("key_drivers", ""), 4)
        rev_insights = _split_text(narratives.get("service_performance", ""), 2)
        cost_insights = _split_text(narratives.get("cost_management", ""), 2)

        # One replacements dict per slide (template has 8 slides)
        slide_replacements = [
            # Slide 1 — Title
            {
                "mbr_period":       period,
                "prepared_by":      "LONGHAUL AI",
                "data_source_count": "3",
            },
            # Slide 2 — Executive Summary
            {
                **kpi_tags,
                "mbr_period":          period,
                "executive_commentary": narratives.get("executive_summary", ""),
            },
            # Slide 3 — Revenue Performance
            {
                "mbr_period":       period,
                "revenue_prior":    kpi_tags["revenue_prior"],
                "total_revenue":    kpi_tags["total_revenue"],
                "revenue_delta":    kpi_tags["revenue_delta"],
                "revenue_insight_1": rev_insights[0],
                "revenue_insight_2": rev_insights[1],
            },
            # Slide 4 — Cost Management
            {
                "mbr_period":         period,
                "cost_per_hour":      kpi_tags["cost_per_hour"],
                "cost_per_hour_delta": kpi_tags["cost_per_hour_delta"],
                "total_cost":         kpi_tags["total_cost"],
                "total_cost_delta":   kpi_tags["total_cost_delta"],
                "cost_insight_1":     cost_insights[0],
                "cost_insight_2":     cost_insights[1],
            },
            # Slide 5 — Operational Efficiency
            {
                "mbr_period":            period,
                "on_time_delivery":      kpi_tags["on_time_delivery"],
                "on_time_delivery_delta": kpi_tags["on_time_delivery_delta"],
                "utilization_rate":      kpi_tags["utilization_rate"],
                "utilization_rate_delta": kpi_tags["utilization_rate_delta"],
                "sla_compliance":        kpi_tags["sla_compliance"],
                "sla_compliance_delta":  kpi_tags["sla_compliance_delta"],
                "efficiency_score":      "",
                "efficiency_analysis":   narratives.get("operational_efficiency", ""),
            },
            # Slide 6 — Sector Performance (mapped to regional data)
            {
                "mbr_period":        period,
                "top_sector_name":   region,
                "top_sector_growth": kpi_tags["revenue_delta"],
                "s1_name":    region,               "s1_revenue": kpi_tags["total_revenue"],
                "s1_growth":  kpi_tags["revenue_delta"], "s1_margin": kpi_tags["operating_margin"],
                "s2_name": "", "s2_revenue": "", "s2_growth": "", "s2_margin": "",
                "s3_name": "", "s3_revenue": "", "s3_growth": "", "s3_margin": "",
                "s4_name": "", "s4_revenue": "", "s4_growth": "", "s4_margin": "",
            },
            # Slide 7 — Key Drivers & Bottom Line
            {
                "mbr_period":         period,
                "driver_1":           drivers[0],
                "driver_2":           drivers[1],
                "driver_3":           drivers[2],
                "driver_4":           drivers[3],
                "bottom_line_summary": narratives.get("bottom_line", ""),
                "call_to_action":     narratives.get("call_to_action", ""),
            },
            # Slide 8 — Data Sources & Methodology
            {
                "mbr_period":        period,
                "data_source_1":     "Fabric Lakehouse — lh_trucking_ops",
                "data_source_2":     "Azure AI Foundry Agent",
                "data_source_3":     "LONGHAUL MBR Platform",
                "data_source_4":     "",
                "data_start_date":   period,
                "data_end_date":     period,
                "report_date":       generated_date,
                "methodology_note":  "KPIs aggregated from the LONGHAUL Fabric Lakehouse via Azure AI Foundry.",
            },
        ]

        with tempfile.TemporaryDirectory() as tmpdir:
            # Download template
            template_bytes = _download_bytes(client, "templates", "mbr_template.pptx")
            template_path = os.path.join(tmpdir, "mbr_template.pptx")
            with open(template_path, "wb") as f:
                f.write(template_bytes)

            # Fill template slide by slide
            prs = Presentation(template_path)
            for idx, slide in enumerate(prs.slides):
                if idx < len(slide_replacements):
                    _fill_slide(slide, slide_replacements[idx])

            # Populate chart shapes (no-op until the Phase 2 template charts exist)
            _fill_charts(prs, analytics, kpis, region)

            # Save filled deck
            deck_filename = f"{region}-{period_slug}-{deck_id}.pptx"
            deck_path = os.path.join(tmpdir, deck_filename)
            prs.save(deck_path)

            # Upload deck to Storage
            with open(deck_path, "rb") as f:
                deck_bytes = f.read()
            deck_blob_path = f"{region}-{period_slug}-{deck_id}.pptx"
            _upload_bytes(client, "decks", deck_blob_path, deck_bytes,
                          "application/vnd.openxmlformats-officedocument.presentationml.presentation")

            # Generate slide thumbnails
            thumb_dir = os.path.join(tmpdir, "thumbs")
            os.makedirs(thumb_dir, exist_ok=True)
            thumbnail_urls: list[str] = []

            try:
                png_paths = _generate_thumbnails(deck_path, thumb_dir)
                for i, png_path in enumerate(png_paths, start=1):
                    blob_path = f"decks/{deck_id}/slide-{i:02d}.png"
                    with open(png_path, "rb") as f:
                        png_bytes = f.read()
                    _upload_bytes(client, "thumbnails", blob_path, png_bytes, "image/png")
                    thumbnail_urls.append(_sas_url(client, "thumbnails", blob_path))
            except Exception as exc:
                logger.warning("Thumbnail generation failed, continuing without thumbnails: %s", exc)

            deck_url = _sas_url(client, "decks", deck_blob_path)

            # Save deck metadata
            metadata = {
                "deck_id":        deck_id,
                "period":         period,
                "region":         region,
                "generated_at":   datetime.now(timezone.utc).isoformat(),
                "deck_blob":      f"decks/{deck_blob_path}",
                "thumbnail_urls": thumbnail_urls,
            }
            _upload_bytes(
                client, "decks-metadata", f"{deck_id}.json",
                json.dumps(metadata).encode(), "application/json",
            )

            logger.info("fill_presentation_template complete: deck_id=%s region=%s period=%s", deck_id, region, period)
            return {"deck_id": deck_id, "deck_url": deck_url, "thumbnail_urls": thumbnail_urls}

    @mcp.tool()
    def get_deck_url(deck_id: str) -> str:
        """Return a short-lived SAS URL for downloading a completed MBR deck.

        Args:
            deck_id: The deck ID returned by fill_presentation_template.

        Returns:
            SAS URL string for the .pptx blob.
        """
        client = _get_blob_client()
        container_client = client.get_container_client("decks")
        for blob in container_client.list_blobs():
            if deck_id in blob.name:
                return _sas_url(client, "decks", blob.name)
        raise ValueError(f"No deck found for deck_id={deck_id!r}")
